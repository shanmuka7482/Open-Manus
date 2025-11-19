from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from app.tool.base import BaseTool, ToolResult
import os
import re

class PptxGenerationTool(BaseTool):
    name: str = "create_presentation"
    description: str = "Creates a PowerPoint presentation with multiple slides, custom fonts, and colors."
    parameters: dict = {
        "type": "object",
        "properties": {
            "filename": {
                "type": "string",
                "description": "The filename for the presentation. Should end with .pptx.",
            },
            "slides": {
                "type": "array",
                "description": "A list of slides to add to the presentation.",
                "items": {
                    "type": "object",
                    "properties": {
                        "title": {"type": "string", "description": "The title of the slide."},
                        "content": {"type": "string", "description": "The content of the slide. Use newlines for bullet points."},
                        "layout": {
                            "type": "string",
                            "description": "Layout type: 'title_slide' (0) or 'title_and_content' (1). Defaults to 'title_and_content'.",
                            "default": "title_and_content"
                        },
                        "font_name": {"type": "string", "description": "Font for the content text. Defaults to 'Arial'."},
                        "font_size": {"type": "integer", "description": "Font size for the content text. Defaults to 18."},
                        "font_color": {"type": "string", "description": "Hex code for font color (e.g., 'FF0000' for red)."},
                        "background_color": {"type": "string", "description": "Hex code for slide background color (e.g., '0000FF' for blue)."}
                    },
                    "required": ["title"]
                }
            }
        },
        "required": ["filename", "slides"],
    }

    async def execute(self, filename: str, slides: list) -> ToolResult:
        """
        Creates an advanced PowerPoint presentation.
        """
        if not filename.endswith(".pptx"):
            return ToolResult(error="Filename must end with .pptx")

        try:
            prs = Presentation()

            # Map layout names to indices
            layout_map = {
                'title_slide': 0,
                'title_and_content': 1,
            }

            for slide_data in slides:
                layout_name = slide_data.get('layout', 'title_and_content').lower()
                layout_index = layout_map.get(layout_name, 1)
                slide_layout = prs.slide_layouts[layout_index]
                slide = prs.slides.add_slide(slide_layout)

                # Set title and content
                if slide.shapes.title:
                    slide.shapes.title.text = slide_data.get('title', '')

                if layout_index == 1 and slide.placeholders[1]:
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.clear() # Clear existing text

                    content_text = slide_data.get('content', '')
                    p = tf.paragraphs[0]
                    p.text = content_text.split('\n')[0] if content_text else ''

                    for line in content_text.split('\n')[1:]:
                        p = tf.add_paragraph()
                        p.text = line
                        p.level = 1

                    # Apply font properties
                    font_name = slide_data.get('font_name', 'Arial')
                    font_size = slide_data.get('font_size', 18)
                    font_color_hex = slide_data.get('font_color')
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.font.name = font_name
                            run.font.size = Pt(font_size)
                            if font_color_hex and re.match(r'^[0-9a-fA-F]{6}$', font_color_hex):
                                run.font.color.rgb = RGBColor.from_string(font_color_hex)

                elif layout_index == 0 and slide.placeholders[1]: # Subtitle for title slide
                    slide.placeholders[1].text = slide_data.get('content', '')


                # Set background color
                bg_color_hex = slide_data.get('background_color')
                if bg_color_hex and re.match(r'^[0-9a-fA-F]{6}$', bg_color_hex):
                    background = slide.background
                    fill = background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(bg_color_hex)

            # Ensure the workspace directory exists
            workspace_dir = os.path.join(os.getcwd(), 'workspace')
            if not os.path.exists(workspace_dir):
                os.makedirs(workspace_dir)

            filepath = os.path.join(workspace_dir, filename)
            prs.save(filepath)

            return ToolResult(output=f"Presentation successfully created at {filepath}")

        except Exception as e:
            return ToolResult(error=f"Failed to create presentation: {str(e)}")
